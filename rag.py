import os
import base64
from pathlib import Path
from typing import Any

import numpy as np
import faiss
from langchain_core.documents import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter
from sentence_transformers import SentenceTransformer
from openai import OpenAI

DEFAULT_EMBEDDING_MODEL = "all-MiniLM-L6-v2"
DEFAULT_LLM_MODEL = "gpt-4.1-mini"
DEFAULT_CHUNK_SIZE = 512
DEFAULT_CHUNK_OVERLAP = 64
DEFAULT_TOP_K = 5

SUPPORTED_EXTENSIONS = {".txt", ".md", ".pdf", ".pptx"}


def _parse_int_setting(name: str, value: Any) -> int:
    try:
        parsed = int(value)
    except (TypeError, ValueError) as exc:
        raise ValueError(f"{name} must be an integer; got {value!r}") from exc
    return parsed


def resolve_config(config: dict[str, Any] | None = None) -> dict[str, Any]:
    """Resolves runtime configuration with defaults and typed settings."""
    config = config or {}

    resolved = {
        "api_key": config.get("api_key", None),
        "base_url": config.get("base_url", None),
        "model": config.get("model", DEFAULT_LLM_MODEL),
        "vision_model": config.get("vision_model") or config.get("model", DEFAULT_LLM_MODEL),
        "embedding_model": config.get("embedding_model", DEFAULT_EMBEDDING_MODEL),
        "top_k": _parse_int_setting(
            "TOP_K", 
            config.get("top_k", DEFAULT_TOP_K)),
        "chunk_size": _parse_int_setting(
            "CHUNK_SIZE", 
            config.get("chunk_size", DEFAULT_CHUNK_SIZE)),
        "chunk_overlap": _parse_int_setting(
            "CHUNK_OVERLAP", 
            config.get("chunk_overlap", DEFAULT_CHUNK_OVERLAP)),
    }

    if resolved["top_k"] <= 0:
        raise ValueError("TOP_K must be > 0")
    if resolved["chunk_size"] <= 0:
        raise ValueError("CHUNK_SIZE must be > 0")
    if resolved["chunk_overlap"] < 0:
        raise ValueError("CHUNK_OVERLAP must be >= 0")
    if resolved["chunk_overlap"] >= resolved["chunk_size"]:
        raise ValueError("CHUNK_OVERLAP must be smaller than CHUNK_SIZE")

    return resolved


def _extract_pdf(file_content: bytes) -> str:
    import fitz
    doc = fitz.open(stream=file_content, filetype="pdf")
    pages = []
    for i, page in enumerate(doc, 1):
        text = page.get_text()
        if text.strip():
            pages.append(f"[Page {i}]\n{text}")
    return "\n\n".join(pages)


def _extract_pptx(file_content: bytes) -> str:
    from io import BytesIO
    from pptx import Presentation
    prs = Presentation(BytesIO(file_content))
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        slide_text = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        slide_text.append(text)
        if slide_text:
            slides.append(f"[Slide {i}]\n" + "\n".join(slide_text))
    return "\n\n".join(slides)


def split_documents(
    docs: list[Document],
    chunk_size: int = DEFAULT_CHUNK_SIZE,
    chunk_overlap: int = DEFAULT_CHUNK_OVERLAP,
) -> list[Document]:
    """Splits documents into overlapping chunks preserving source metadata."""
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=chunk_size,
        chunk_overlap=chunk_overlap,
    )
    return splitter.split_documents(docs)


def build_index(
    chunks: list[Document],
    embedding_model: SentenceTransformer,
) -> faiss.IndexFlatIP:
    """Creates a FAISS inner-product index for embedded document chunks."""
    texts = [chunk.page_content for chunk in chunks]
    embeddings = embedding_model.encode(texts, convert_to_numpy=True).astype("float32")
    faiss.normalize_L2(embeddings)
    index = faiss.IndexFlatIP(embeddings.shape[1])
    index.add(embeddings)
    return index


def retrieve(
    query: str,
    index: faiss.IndexFlatIP,
    model: SentenceTransformer,
    chunks: list[Document],
    k: int = DEFAULT_TOP_K,
) -> list[dict]:
    """Gets the most relevant chunks for a query, ordered by similarity."""
    query_embedding = model.encode([query], convert_to_numpy=True).astype("float32")
    faiss.normalize_L2(query_embedding)

    distances, indices = index.search(query_embedding, k)
    results = []
    
    for idx, score in zip(indices[0], distances[0]):
        chunk = chunks[idx]
        results.append({
            "text": chunk.page_content,
            "score": float(score),
            "metadata": chunk.metadata,
        })
    return results


SYSTEM_PROMPT = (
    """You are StudyBot, an intelligent study assistant. 
    You help students learn and understand their uploaded study materials. 
    You have access to the student's materials (PDFs, presentations, notes, images) 
    Your capabilities include: answering questions accurately from the context, 
    generating clear structured summaries, creating quiz questions to test understanding, and explaining complex concepts in simpler terms. 
    Always base your responses strictly on the retrieved context. 
    If a question cannot be answered from the available materials, say so clearly. 
    Mention the source filename when it is helpful to the student."""
)


class Assistant:
    """Stateful RAG study assistant.

    The assistant owns the pipeline components, resolved configuration, and
    conversation history. It supports Q&A, summarization, and quiz generation
    from user-uploaded study materials.
    """

    def __init__(
        self,
        index: faiss.IndexFlatIP,
        model: SentenceTransformer,
        chunks: list[Document],
        client: OpenAI,
        docs: list[Document],
        config: dict[str, Any] | None = None,
    ) -> None:
        self.index = index
        self.model = model
        self.chunks = chunks
        self.client = client
        self.docs = docs
        self.config = resolve_config(config)
        self.llm_model = self.config["model"]
        self.top_k = self.config["top_k"]
        self.history: list[dict[str, str]] = []

    def ask(self, question: str, k: int | None = None) -> str:
        """Generates an answer from retrieved context and conversation history."""
        if not self.chunks:
            return (
                "No study materials are loaded. "
                "Please add files to the uploads/ folder and restart."
            )

        if k is None:
            k = self.top_k

        retrieved = retrieve(question, self.index, self.model, self.chunks, k=k)

        if retrieved and retrieved[0]["score"] == 0:
            return "No relevant information found in the uploaded materials."

        context = "\n\n".join([
            f"[Source: {r['metadata']['filename']} ({r['metadata']['type']})]\n{r['text']}"
            for r in retrieved
        ])

        context_question = f"<Context>{context}</Context><Question>{question}</Question>"

        response = self.client.chat.completions.create(
            model=self.llm_model,
            messages=[
                {"role": "system", "content": SYSTEM_PROMPT},
                *self.history,
                {"role": "user", "content": context_question},
            ],
        ).choices[0].message.content.strip()

        self.history.append({"role": "user", "content": context_question})
        self.history.append({"role": "assistant", "content": response})

        return response

    def summarize(self, topic: str | None = None) -> dict[str, Any]:
        """Generates a structured study summary, optionally focused on a topic.
        
        Returns the response and a list of sources used for the summary 
        """
        if not self.chunks:
            return {
                "response": "No study materials are loaded.",
                "sources": []
            }

        query = "overview main topics key concepts" if topic is None else f"{topic} summary overview"
        k = min(self.top_k * 2, len(self.chunks))
        retrieved = retrieve(query, self.index, self.model, self.chunks, k=k)

        if not retrieved or retrieved[0]["score"] == 0:
            return {
                "response": "No relevant material found to summarize.",
                "sources": []
            }

        context = "\n\n".join([
            f"[Source: {r['metadata']['filename']}]\n{r['text']}"
            for r in retrieved
        ])
        prompt = (
            "Generate a comprehensive, well-structured study summary"
            + (f" focused on the topic: '{topic}'" if topic else " of the provided materials")
            + ". Include key concepts, definitions, and important points organized clearly."
        )

        response = self.client.chat.completions.create(
            model=self.llm_model,
            messages=[
                {"role": "system", "content": SYSTEM_PROMPT},
                {"role": "user", "content": f"<Context>{context}</Context><Question>{prompt}</Question>"},
            ],
        ).choices[0].message.content.strip()

        self.history.append({"role": "user", "content": f"<Context>{context}</Context><Question>{prompt}</Question>"})
        self.history.append({"role": "assistant", "content": response})

        sources = []
        seen = set()
        for r in retrieved:
            filename = r['metadata']['filename']
            if filename not in seen:
                seen.add(filename)
                sources.append({
                    "filename": filename,
                    "type": r['metadata']['type'],
                    "score": r['score']
                })

        return {
            "response": response,
            "sources": sources
        }

    def generate_quiz(self, topic: str | None = None, num_questions: int = 5) -> dict[str, Any]:
        """Generates multiple-choice quiz questions from the study materials.
        
        Returns the generated quiz and a list of sources used for question generation.
        """
        if not self.chunks:
            return {
                "response": "No study materials are loaded.",
                "sources": []
            }

        query = topic or "key concepts important facts definitions"
        k = min(self.top_k * 2, len(self.chunks))
        retrieved = retrieve(query, self.index, self.model, self.chunks, k=k)

        if not retrieved or retrieved[0]["score"] == 0:
            return {
                "response": "No relevant material found to generate quiz questions.",
                "sources": []
            }

        context = "\n\n".join([
            f"[Source: {r['metadata']['filename']}]\n{r['text']}"
            for r in retrieved
        ])
        prompt = (
            f"Generate {num_questions} quiz questions"
            + (f" about '{topic}'" if topic else " based on the study materials")
            + ". For each question provide 4 multiple-choice options (A, B, C, D) "
            + "and indicate the correct answer with a brief explanation."
        )

        response = self.client.chat.completions.create(
            model=self.llm_model,
            messages=[
                {"role": "system", "content": SYSTEM_PROMPT},
                {"role": "user", "content": f"<Context>{context}</Context><Question>{prompt}</Question>"},
            ],
        ).choices[0].message.content.strip()

        self.history.append({"role": "user", "content": f"<Context>{context}</Context><Question>{prompt}</Question>"})
        self.history.append({"role": "assistant", "content": response})

        # Extract unique sources
        sources = []
        seen = set()
        for r in retrieved:
            filename = r['metadata']['filename']
            if filename not in seen:
                seen.add(filename)
                sources.append({
                    "filename": filename,
                    "type": r['metadata']['type'],
                    "score": r['score']
                })

        return {
            "response": response,
            "sources": sources
        }

    def list_files(self) -> str:
        """Returns a formatted list of all loaded study materials."""
        if not self.docs:
            return "No files loaded. Add study materials to the uploads/ folder and restart."
        seen: set[str] = set()
        lines = []
        for doc in self.docs:
            name = doc.metadata["filename"]
            if name not in seen:
                seen.add(name)
                lines.append(f"  [{doc.metadata['type']}] {name}")
        return "Loaded study materials:\n" + "\n".join(lines)

    def list_files_json(self) -> list[dict]:
        """Returns loaded files as structured dicts for JSON serialization."""
        seen: set[str] = set()
        files = []
        for doc in self.docs:
            name = doc.metadata["filename"]
            if name not in seen:
                seen.add(name)
                files.append({
                    "filename": name,
                    "type": doc.metadata["type"],
                    "source": doc.metadata["source"],
                })
        return files

    def clear(self) -> None:
        """Resets all documents, chunks, history, and FAISS index."""
        self.docs = []
        self.chunks = []
        self.history = []
        dim = self.model.get_sentence_embedding_dimension()
        self.index = faiss.IndexFlatIP(dim)

    def remove_file(self, filename: str) -> bool:
        """Removes a file by filename and rebuilds the index."""
        # Select all documents except the one that matches the filename to remove
        initial_count = len(self.docs)
        self.docs = [doc for doc in self.docs if doc.metadata["filename"] != filename]
        if len(self.docs) == initial_count:
            return False
        self.history = []
        
        # Remake chunks and index
        self.chunks = split_documents(
            self.docs,
            chunk_size=self.config["chunk_size"],
            chunk_overlap=self.config["chunk_overlap"],
        )
        if self.chunks:
            self.index = build_index(self.chunks, self.model)
        else:
            dim = self.model.get_sentence_embedding_dimension()
            self.index = faiss.IndexFlatIP(dim)
        return True

    def add_documents_from_content(self, files: list[tuple[str, bytes]]) -> dict[str, bool]:
        """Adds documents from file content (filename, bytes) and rebuilds the index.
        
        Returns a dictionary mapping filenames to boolean values indicating success.
        """
        results = {}
        new_docs = []
        
        for filename, file_content in files:
            if any(doc.metadata["filename"] == filename for doc in self.docs):
                results[filename] = False
                continue

            ext = Path(filename).suffix.lower()
            if ext not in SUPPORTED_EXTENSIONS:
                results[filename] = False
                continue
            
            try:
                if ext in {".txt", ".md"}:
                    content = file_content.decode(encoding="utf-8", errors="replace")
                    doc_type = "text"
                elif ext == ".pdf":
                    content = _extract_pdf(file_content)
                    doc_type = "pdf"
                elif ext == ".pptx":
                    content = _extract_pptx(file_content)
                    doc_type = "presentation"
                else:
                    results[filename] = False
                    continue
                
                if content.strip():
                    metadata = {
                        "source": filename,
                        "type": doc_type,
                        "filename": filename,
                    }
                    doc = Document(page_content=content, metadata=metadata)
                    new_docs.append(doc)
                    results[filename] = True
                else:
                    results[filename] = False
            except Exception as e:
                print(f"Error: could not load {filename} for upload: {e}")
                results[filename] = False
        
        # If successful, add new documents and rebuild chunks and index
        if new_docs:
            self.docs.extend(new_docs)
            
            self.chunks = split_documents(
                self.docs,
                chunk_size=self.config["chunk_size"],
                chunk_overlap=self.config["chunk_overlap"],
            )
            self.index = build_index(self.chunks, self.model)
        
        return results

    @classmethod
    def from_config(cls, config: dict[str, Any] | None = None) -> "Assistant":
        """Initializes the full pipeline and returns a ready Assistant.

        Creates the API client first so image files can be processed via the
        vision API during document loading.
        """
        resolved_config = resolve_config(config)

        print("Initializing...")
        docs = []

        print("Splitting into chunks...")
        chunks = split_documents(
            docs, 
            chunk_size=resolved_config["chunk_size"], 
            chunk_overlap=resolved_config["chunk_overlap"]
        )
        print(f"  Created {len(chunks)} chunks")

        embedding_model = SentenceTransformer(resolved_config["embedding_model"])

        print("Building FAISS index...")
        if chunks:
            index = build_index(chunks, embedding_model)
            print(f"  Indexed {index.ntotal} vectors (dim={index.d})")
        else:
            dim = embedding_model.get_sentence_embedding_dimension()
            index = faiss.IndexFlatIP(dim)
            print("(Empty index, no documents loaded)")

        client_kwargs = {}
        if resolved_config["api_key"]:
            client_kwargs["api_key"] = resolved_config["api_key"]
        if resolved_config["base_url"]:
            client_kwargs["base_url"] = resolved_config["base_url"]
        client = OpenAI(**client_kwargs)

        print("Ready!\n")
        return cls(index, embedding_model, chunks, client, docs, resolved_config)
